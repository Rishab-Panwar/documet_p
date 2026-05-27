from __future__ import annotations
from pathlib import Path
from typing import Iterable, List
from fastapi import UploadFile
from langchain.schema import Document
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from logger import GLOBAL_LOGGER as log
from exception.custom_exception import DocumentPortalException
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx", ".md", ".csv", ".xlsx", ".xls", ".db", ".sqlite", ".sqlite3"}


def load_documents(paths: Iterable[Path]) -> List[Document]:
    """Load docs using appropriate loader based on extension."""
    docs: List[Document] = []
    try:
        for p in paths:
            ext = p.suffix.lower()
            try:
                if ext == ".pdf":
                    docs.extend(PyPDFLoader(str(p)).load())
                elif ext == ".docx":
                    docs.extend(Docx2txtLoader(str(p)).load())
                elif ext == ".txt":
                    docs.extend(TextLoader(str(p), encoding="utf-8").load())
                elif ext == ".pptx":
                    docs.append(_doc_from_text(_read_pptx(p), p, {"file_type": ext}))
                elif ext == ".md":
                    docs.append(_doc_from_text(p.read_text(encoding="utf-8"), p, {"file_type": ext}))
                elif ext == ".csv":
                    docs.append(_doc_from_text(_read_csv(p), p, {"file_type": ext}))
                elif ext == ".xlsx":
                    for sheet_name, text in _read_xlsx_as_sheets(p):
                        docs.append(_doc_from_text(text, p, {"file_type": ext, "sheet": sheet_name}))
                elif ext == ".xls":
                    for sheet_name, text in _read_xls_as_sheets(p):
                        docs.append(_doc_from_text(text, p, {"file_type": ext, "sheet": sheet_name}))
                elif ext in {".db", ".sqlite", ".sqlite3"}:
                    docs.append(_doc_from_text(_read_sqlite_dump(p), p, {"file_type": ext}))
                else:
                    log.warning("Unsupported extension skipped", path=str(p))
            except Exception as e:
                log.error("Failed loading a document", path=str(p), error=str(e))
                # continue with next file
        log.info("Documents loaded", count=len(docs))
        return docs
    except Exception as e:
        log.error("Failed loading documents", error=str(e))
        raise DocumentPortalException("Error loading documents", e) from e

def concat_for_analysis(docs: List[Document]) -> str:
    parts = []
    for d in docs:
        src = d.metadata.get("source") or d.metadata.get("file_path") or "unknown"
        parts.append(f"\n--- SOURCE: {src} ---\n{d.page_content}")
    return "\n".join(parts)

def concat_for_comparison(ref_docs: List[Document], act_docs: List[Document]) -> str:
    left = concat_for_analysis(ref_docs)
    right = concat_for_analysis(act_docs)
    return f"<<REFERENCE_DOCUMENTS>>\n{left}\n\n<<ACTUAL_DOCUMENTS>>\n{right}"

# ---------- Helpers ----------
class FastAPIFileAdapter:
    """Adapt FastAPI UploadFile -> .name + .getbuffer() API"""
    def __init__(self, uf: UploadFile):
        self._uf = uf
        self.name = uf.filename
    def getbuffer(self) -> bytes:
        self._uf.file.seek(0)
        return self._uf.file.read()

def read_pdf_via_handler(handler, path: str) -> str:
    if hasattr(handler, "read_pdf"):
        return handler.read_pdf(path)  # type: ignore
    if hasattr(handler, "read_"):
        return handler.read_(path)  # type: ignore
    raise RuntimeError("DocHandler has neither read_pdf nor read_ method.")


def _doc_from_text(text: str, p: Path, extra_md: dict) -> Document:
    md = {"source": str(p.resolve())}
    if extra_md:
        md.update(extra_md)
    return Document(page_content=text or "", metadata=md)


def _read_pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            parts.append(f"\n--- Slide {i} ---\n" + "\n".join(slide_text))
    return "\n".join(parts)


def _read_csv(p: Path) -> str:
    import csv
    lines: List[str] = []
    with p.open("r", encoding="utf-8", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            lines.append(", ".join("" if c is None else str(c) for c in row))
    return "\n".join(lines)


def _read_xlsx_as_sheets(p: Path):
    import openpyxl
    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        parts = [f"\n--- Sheet: {ws.title} ---"]
        for row in ws.iter_rows(values_only=True):
            parts.append("\t".join("" if c is None else str(c) for c in row))
        out.append((ws.title, "\n".join(parts)))
    return out


def _read_xls_as_sheets(p: Path):
    import xlrd
    wb = xlrd.open_workbook(str(p))
    out = []
    for sheet in wb.sheets():
        parts = [f"\n--- Sheet: {sheet.name} ---"]
        for rx in range(sheet.nrows):
            row = [sheet.cell_value(rx, cx) for cx in range(sheet.ncols)]
            parts.append("\t".join("" if c is None else str(c) for c in row))
        out.append((sheet.name, "\n".join(parts)))
    return out


def _read_sqlite_dump(p: Path) -> str:
    import sqlite3
    uri = f"file:{p}?mode=ro"
    conn = sqlite3.connect(uri, uri=True)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()

    cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%' ORDER BY name;")
    tables = [r[0] for r in cur.fetchall()]
    if not tables:
        conn.close()
        return ""

    parts: List[str] = []
    for t in tables:
        parts.append(f"\n--- Table: {t} ---")
        try:
            cur.execute(f"PRAGMA table_info('{t}')")
            cols = [row[1] for row in cur.fetchall()]
            if cols:
                parts.append("# Columns: " + ", ".join(cols))
        except Exception:
            pass
        try:
            cur.execute(f"SELECT * FROM '{t}' LIMIT 1000")
            rows = cur.fetchall()
            for r in rows:
                vals = [r[k] for k in r.keys()] if isinstance(r, sqlite3.Row) else list(r)
                parts.append("\t".join("" if v is None else str(v) for v in vals))
        except Exception as e:
            parts.append(f"# Error reading table '{t}': {e}")

    conn.close()
    return "\n".join(parts)
